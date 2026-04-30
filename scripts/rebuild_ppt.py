#!/usr/bin/env python3
"""
精简PPT到18页 - 学术报告版
- 重做第5页（AI三阶段流程）
- 删除：合规边界、发展路线图、未来展望
- 合并：真实图像展示页
- 新增：奖项页
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import shutil

INPUT_FILE = '/tmp/管院长-正式稿.pptx'
OUTPUT_FILE = '/root/.openclaw/workspace/attachments/管院长-正式稿_学术版.pptx'
ASSETS_DIR = '/root/.openclaw/workspace/attachments/ppt_assets'

prs = Presentation(INPUT_FILE)
total_orig = len(prs.slides)

# 配色
BLUE = RGBColor(0x1a, 0x73, 0xe8)
BLUE_DARK = RGBColor(0x0d, 0x47, 0xa1)
GREEN = RGBColor(0x34, 0xa8, 0x53)
RED = RGBColor(0xea, 0x43, 0x35)
ORANGE = RGBColor(0xff, 0x6d, 0x00)
PURPLE = RGBColor(0x9c, 0x27, 0xb0)
TEAL = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x20, 0x21, 0x24)
GRAY = RGBColor(0x5f, 0x63, 0x68)
GRAY_LIGHT = RGBColor(0xe8, 0xea, 0xed)
GOLD = RGBColor(0xff, 0xb3, 0x00)
SILVER = RGBColor(0x90, 0xa4, 0xae)
BRONZE = RGBColor(0xff, 0x8a, 0x65)

def set_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return txBox

def add_header(slide, title, subtitle, num, total):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.6), fill_color=BLUE)
    add_text(slide, Inches(0.8), Inches(0.05), Inches(10), Inches(0.5), title, size=22, bold=True, color=WHITE)
    add_text(slide, Inches(0.8), Inches(0.65), Inches(10), Inches(0.35), subtitle, size=13, color=GRAY)
    add_shape(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), fill_color=RGBColor(0xf8, 0xf9, 0xfa))
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.3), 'AIMED 充盈视界 · 超声造影人工智能联合实验室（第一医院体检中心+南京大学）', size=9, color=GRAY)
    add_text(slide, Inches(12), Inches(7.12), Inches(1), Inches(0.3), f'{num} / {total}', size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# 18页结构规划
# ============================================================
# 1  封面
# 2  目录
# 3  临床背景与痛点
# 4  口服超声充盈造影技术
# 5  AI三阶段智能诊断流程（重做 - 学术风格）
# 6  AI系统架构全景图
# 7  Su-RADS 胰腺分类标准
# 8  充盈造影对比图（模拟）
# 9  真实超声图像展示（合并原9-12）
# 10 真实病例对比 - 充盈前后
# 11 充盈超声造影 vs 胃镜/常规超声
# 12 标准文档真实超声图像
# 13 临床应用案例演示
# 14 核心优势总结
# 15 技术架构概览
# 16 AI技术原理
# 17 项目荣誉与奖项（新增）
# 18 感谢聆听 + Q&A


# ============================================================
# 创建新PPT
# ============================================================
prs_new = Presentation()
prs_new.slide_width = Inches(13.333)
prs_new.slide_height = Inches(7.5)
TOTAL = 18


# ============================================================
# Slide 1: 封面（复制原Slide 1）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=BLUE)
add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), '口服超声充盈造影', size=36, bold=True, color=BLUE)
add_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1), '对胰腺疾病诊断的 AI 模型', size=36, bold=True, color=BLUE)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6), '基于《胃癌超声初筛临床应用中国专家共识意见（2025年版）》诊断逻辑延伸', size=14, color=GRAY)
add_shape(slide, Inches(1.5), Inches(4.2), Inches(3), Inches(0.04), fill_color=BLUE)
add_text(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.4), '2026 医学人工智能前沿学术会议', size=16, color=BLACK)
add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4), '湖州第一医院 · 体检中心', size=14, color=GRAY)
add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4), '演讲人：管建明 院长', size=12, color=GRAY)
add_shape(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), fill_color=BLUE)
print("✅ Slide 1: 封面")


# ============================================================
# Slide 2: 目录（复制原Slide 2）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '汇报目录', 'Agenda', 2, TOTAL)

agenda_items = [
    ('01', '临床背景与痛点', '胰腺疾病诊断现状与挑战', BLUE),
    ('02', '口服超声充盈造影', '技术原理与胰腺成像优势', GREEN),
    ('03', 'AI 模型架构设计', '三阶段诊断流程', BLUE),
    ('04', 'Su-RADS 胰腺分类', 'AI 辅助分级标准', ORANGE),
    ('05', '充盈造影对比展示', '疾病图像与 AI 标注对比', TEAL),
    ('06', '临床应用案例', '模拟演示与效果展示', PURPLE),
    ('07', '项目荣誉与奖项', '科研团队成果展示', GOLD),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(1.2) + Inches(i * 0.8)
    add_rounded(slide, Inches(1.2), y, Inches(0.6), Inches(0.6), fill_color=color)
    add_text(slide, Inches(1.2), y + Inches(0.08), Inches(0.6), Inches(0.45), num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.0), y + Inches(0.02), Inches(5), Inches(0.35), title, size=16, bold=True, color=BLACK)
    add_text(slide, Inches(2.0), y + Inches(0.38), Inches(8), Inches(0.3), desc, size=11, color=GRAY)
print("✅ Slide 2: 目录")


# ============================================================
# Slide 3: 临床背景（复制原Slide 3）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '01 临床背景', '胰腺疾病诊断现状与挑战', 3, TOTAL)

add_rounded(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), fill_color=RGBColor(0xe3, 0xf2, 0xfd), line_color=BLUE, line_width=Pt(1))
add_text(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5), '📊 胰腺疾病流行病学', size=16, bold=True, color=BLUE)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['中国胰腺癌发病率逐年上升', '5 年生存率仅约 10%', '早期诊断率不足 20%', '体检人群中胰腺异常检出率约 3-5%', '中老年人群是高危群体']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.name = '微软雅黑'
print("✅ Slide 3: 临床背景")


# ============================================================
# Slide 4: 充盈造影技术（复制原Slide 4）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '02 口服超声充盈造影技术', '技术原理与胰腺成像优势', 4, TOTAL)

add_rounded(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(4.5), fill_color=WHITE, line_color=BLUE, line_width=Pt(1))
add_text(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5), '🔬 技术原理', size=16, bold=True, color=BLUE)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.2), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['患者口服专用超声造影剂', '造影剂在胃肠道内形成均匀充盈', '消除胃肠道气体干扰', '胃腔作为「声学窗」', '透过胃窗观察胰腺等邻近器官', '获得清晰、稳定的胰腺图像']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.name = '微软雅黑'

add_rounded(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(4.5), fill_color=WHITE, line_color=GREEN, line_width=Pt(1))
add_text(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5), '✅ 核心优势', size=16, bold=True, color=GREEN)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.6), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['🟢 无创无辐射，患者接受度高', '🟢 操作简单，适合体检筛查', '🟢 成本低廉，可重复检查', '🟢 成像质量稳定，减少气体干扰', '🟢 可动态观察胰腺血流情况', '🟢 适合基层医院推广']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.name = '微软雅黑'

add_rounded(slide, Inches(0.5), Inches(6.0), Inches(12.5), Inches(0.8), fill_color=RGBColor(0xe8, 0xf5, 0xe9), line_color=GREEN, line_width=Pt(1))
add_text(slide, Inches(0.8), Inches(6.15), Inches(12), Inches(0.5), '💡 技术依据：《胃癌超声初筛临床应用中国专家共识意见（2025年版）》已明确口服超声充盈造影在胃部疾病筛查中的价值和标准化流程', size=11, color=BLUE_DARK)
print("✅ Slide 4: 充盈造影技术")


# ============================================================
# Slide 5: AI三阶段智能诊断流程（重做 - 学术风格）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '03 AI 模型架构设计', '三阶段智能诊断流程', 5, TOTAL)

# 插入重新生成的AI三阶段图
stage_img = os.path.join(ASSETS_DIR, 'ai_3stage_diagram.png')
if os.path.exists(stage_img):
    slide.shapes.add_picture(stage_img, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8))
print("✅ Slide 5: AI三阶段智能诊断流程（重做）")


# ============================================================
# Slide 6: AI系统架构全景图（复制原Slide 6）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '03 AI 模型架构设计', '系统架构全景图', 6, TOTAL)

arch_img = os.path.join(ASSETS_DIR, 'ai_system_architecture.png')
if os.path.exists(arch_img):
    slide.shapes.add_picture(arch_img, Inches(0.2), Inches(1.0), Inches(12.9), Inches(5.9))
print("✅ Slide 6: AI系统架构全景图")


# ============================================================
# Slide 7: Su-RADS 胰腺分类标准（复制原Slide 7）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '04 Su-RADS 胰腺分类标准', 'AI 辅助分级体系', 7, TOTAL)

surads_img = os.path.join(ASSETS_DIR, 'surads_chart.png')
if os.path.exists(surads_img):
    slide.shapes.add_picture(surads_img, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))
print("✅ Slide 7: Su-RADS 胰腺分类标准")


# ============================================================
# Slide 8: 充盈造影对比图（复制原Slide 8）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '05 充盈造影对比展示', '疾病充盈造影前后图像对比 + AI 标注', 8, TOTAL)

contrast_img = os.path.join(ASSETS_DIR, 'contrast_comparison.png')
if os.path.exists(contrast_img):
    slide.shapes.add_picture(contrast_img, Inches(0.2), Inches(1.0), Inches(12.9), Inches(5.8))
print("✅ Slide 8: 充盈造影对比图")


# ============================================================
# Slide 9: 真实超声图像展示（合并原9-12为一页）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '05 充盈造影对比展示', '真实病例超声图像（充盈后）', 9, TOTAL)

# 选取6张真实超声图像 3x2 布局
real_images = [
    ('data/cases/stomach/2026-04-23/胃底部恶性间质瘤_A_g_20260423011923.jpg', '胃底部恶性间质瘤'),
    ('data/cases/stomach/2026-04-23/胃角溃疡（活动期）A_o_20260423012114.jpg', '胃角溃疡（活动期）'),
    ('data/cases/stomach/2026-04-23/胃小弯粘膜下良性间质瘤_L_20260423010552.jpg', '胃小弯良性间质瘤'),
    ('data/cases/stomach/2026-04-23/胃体部早期胃 MT 凹陷型 1_1776906580.jpg', '胃体部早期 MT 凹陷型'),
    ('data/cases/stomach/2026-04-23/胃恶性淋巴瘤_aq_20260423012731.jpg', '胃恶性淋巴瘤'),
    ('data/cases/pancreas/2026-04-23/case_20260423000323_62e6efa2.jpg', '胰腺病例'),
]

for i, (img_path, label) in enumerate(real_images):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.2) + Inches(row * 2.8)
    
    full_path = os.path.join('/root/.openclaw/workspace', img_path)
    if os.path.exists(full_path):
        slide.shapes.add_picture(full_path, x, y, Inches(3.8), Inches(2.3))
    
    # 标签
    add_rounded(slide, x, y + Inches(2.3), Inches(3.8), Inches(0.35), fill_color=BLUE)
    add_text(slide, x, y + Inches(2.32), Inches(3.8), Inches(0.3), label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("✅ Slide 9: 真实超声图像展示")


# ============================================================
# Slide 10: 真实病例对比（复制原Slide 10）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '05 充盈造影对比展示', '真实病例对比 - 充盈前后', 10, TOTAL)

comparison_groups = [
    [
        ('data/cases/stomach/2026-04-23/胃底部恶性间质瘤（粘膜下型）_c_20260423010519.jpg', '充盈前'),
        ('data/cases/stomach/2026-04-23/胃底部恶性间质瘤_A_g_20260423011923.jpg', '充盈后'),
    ],
    [
        ('data/cases/stomach/2026-04-23/胃角溃疡_R_r_20260423012154.jpg', '充盈前'),
        ('data/cases/stomach/2026-04-23/胃角溃疡（活动期）A_o_20260423012114.jpg', '充盈后'),
    ],
    [
        ('data/cases/stomach/2026-04-23/胃小弯粘膜下良性间质瘤_i_20260423012029.jpg', '充盈前'),
        ('data/cases/stomach/2026-04-23/胃小弯粘膜下良性间质瘤_L_20260423010552.jpg', '充盈后'),
    ],
]

for i, group in enumerate(comparison_groups):
    x = Inches(0.5) + Inches(i * 4.2)
    y = Inches(1.3)
    
    for j, (img_path, label) in enumerate(group):
        full_path = os.path.join('/root/.openclaw/workspace', img_path)
        if os.path.exists(full_path):
            slide.shapes.add_picture(full_path, x, y + Inches(j * 2.5), Inches(3.8), Inches(2.2))
        
        # 标签
        label_color = GRAY if label == '充盈前' else GREEN
        add_rounded(slide, x, y + Inches(2.2) + Inches(j * 2.5), Inches(3.8), Inches(0.3), fill_color=label_color)
        add_text(slide, x, y + Inches(2.22) + Inches(j * 2.5), Inches(3.8), Inches(0.28), label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("✅ Slide 10: 真实病例对比")


# ============================================================
# Slide 11: 文档真实对比图（复制原Slide 11）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '05 充盈造影对比展示', '充盈超声造影 vs 胃镜/常规超声 — 真实对比', 11, TOTAL)

doc_images = [
    ('/tmp/doc_images/image_0.jpeg', '充盈超声造影 — 胃壁全层清晰显示'),
    ('/tmp/doc_images/image_3.jpeg', '胃镜 — 仅显示黏膜表面'),
    ('/tmp/doc_images/image_12.jpeg', '充盈超声 — 胰腺结构完整显示'),
    ('/tmp/doc_images/image_10.jpeg', '常规超声 — 气体干扰胰腺显示不清'),
]

for i, (img_path, label) in enumerate(doc_images):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.2) + Inches(row * 2.9)
    
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, Inches(5.8), Inches(2.5))
    
    # 标签
    label_color = GREEN if '充盈' in label else GRAY
    add_rounded(slide, x, y + Inches(2.5), Inches(5.8), Inches(0.35), fill_color=label_color)
    add_text(slide, x, y + Inches(2.52), Inches(5.8), Inches(0.3), label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4), '💡 来源：《充盈超声造影与胃镜/常规超声对比：理想的普筛早筛工具》', size=10, color=GRAY)
print("✅ Slide 11: 文档真实对比图")


# ============================================================
# Slide 12: 标准文档真实超声图像（复制原Slide 12）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '05 充盈造影对比展示', '标准文档真实超声图像 — 管建明修订稿', 12, TOTAL)

standard_images = [
    ('/tmp/doc2_images/img_1.png', '造影前空腹 vs 造影后充盈对比'),
    ('/tmp/doc2_images/img_5.png', '胰头标准切面（充盈造影后）'),
    ('/tmp/doc2_images/img_6.png', '胰体标准切面（充盈造影后）'),
    ('/tmp/doc2_images/img_13.png', '胰尾长轴切面（充盈造影后）'),
    ('/tmp/doc2_images/img_15.png', '门静脉胆总管切面'),
    ('/tmp/doc2_images/img_17.png', '腹主动脉胰体切面'),
]

for i, (img_path, label) in enumerate(standard_images):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.2) + Inches(row * 2.9)
    
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, Inches(3.8), Inches(2.5))
    
    add_rounded(slide, x, y + Inches(2.5), Inches(3.8), Inches(0.35), fill_color=BLUE)
    add_text(slide, x, y + Inches(2.52), Inches(3.8), Inches(0.3), label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4), '💡 来源：《标准胰腺超声造影诊断检查标准和 AI 数据采集标准》（管建明修订稿，2025年10月）', size=10, color=GRAY)
print("✅ Slide 12: 标准文档真实超声图像")


# ============================================================
# Slide 13: 临床应用案例（复制原Slide 13）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '06 临床应用案例演示', '模拟案例展示 AI 辅助诊断流程', 13, TOTAL)

add_rounded(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), fill_color=WHITE, line_color=BLUE, line_width=Pt(1))
add_text(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5), '📋 案例一：中年男性，体检筛查', size=15, bold=True, color=BLUE)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['患者信息：', '男，52 岁，无明显症状', '常规体检，有吸烟史', '', '检查过程：', '1. 口服造影剂 500ml', '2. 等待 15 分钟充盈', '3. 超声检查胰腺', '', 'AI 分析结果：', 'Su-RADS 1 · 阴性', '胰腺形态正常，无异常发现']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = BLACK if not line.startswith('Su') else GREEN
    p.font.bold = line.endswith('：') or line.startswith('Su')
    p.font.name = '微软雅黑'

add_rounded(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5), fill_color=WHITE, line_color=ORANGE, line_width=Pt(1))
add_text(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5), '📋 案例二：老年女性，腹痛就诊', size=15, bold=True, color=ORANGE)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['患者信息：', '女，65 岁，上腹痛 2 周', '体重下降，食欲减退', '', '检查过程：', '1. 口服造影剂 500ml', '2. 等待 15 分钟充盈', '3. 超声检查胰腺', '', 'AI 分析结果：', 'Su-RADS 4 · 高危险度', '可疑恶性，建议增强 CT/MRI']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = BLACK if not line.startswith('Su') else RED
    p.font.bold = line.endswith('：') or line.startswith('Su')
    p.font.name = '微软雅黑'
print("✅ Slide 13: 临床应用案例")


# ============================================================
# Slide 14: 核心优势总结（复制原Slide 16）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '核心优势总结', '为什么选择口服超声充盈造影 + AI', 14, TOTAL)

adv_img = os.path.join(ASSETS_DIR, 'advantages_chart.png')
if os.path.exists(adv_img):
    slide.shapes.add_picture(adv_img, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))
print("✅ Slide 14: 核心优势总结")


# ============================================================
# Slide 15: 技术架构概览（复制原Slide 17）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '技术架构概览', '系统如何工作（非技术视角）', 15, TOTAL)

layers = [
    ('🖥️ 用户界面层', '医生/患者操作界面', BLUE, Inches(1.3)),
    ('🧠 AI 分析层', '图像识别 + 智能诊断', GREEN, Inches(2.7)),
    ('📊 数据层', '病例数据库 + 知识库', ORANGE, Inches(4.1)),
    ('🔐 存证层', '区块链存证 + 审计', PURPLE, Inches(5.5)),
]

for title, desc, color, y in layers:
    add_rounded(slide, Inches(1.5), y, Inches(10), Inches(1.1), fill_color=WHITE, line_color=color, line_width=Pt(2))
    add_text(slide, Inches(2.0), y + Inches(0.15), Inches(4), Inches(0.4), title, size=16, bold=True, color=color)
    add_text(slide, Inches(6.0), y + Inches(0.35), Inches(5), Inches(0.4), desc, size=13, color=GRAY)
    if color != PURPLE:
        add_text(slide, Inches(11.5), y + Inches(0.3), Inches(0.5), Inches(0.5), '↓', size=20, color=color, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4), '💡 所有数据本地化部署，符合医疗数据安全要求', size=11, color=GRAY)
print("✅ Slide 15: 技术架构概览")


# ============================================================
# Slide 16: AI技术原理（复制原Slide 18）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, 'AI 技术原理（基础入门）', '无需编程知识，理解 AI 如何辅助诊断', 16, TOTAL)

add_rounded(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), fill_color=WHITE, line_color=BLUE, line_width=Pt(1))
add_text(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5), '🤖 AI 在医学影像中的角色', size=16, bold=True, color=BLUE)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['AI = 计算机视觉 + 深度学习', '', '类比理解：', '就像培养一个「超级助手」', '看过 10000+ 张超声图像后', '能自动识别可疑病灶', '', 'AI 不是替代医生，而是：', '✓ 提高诊断效率', '✓ 减少漏诊误诊', '✓ 统一诊断标准']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = GREEN if line.startswith('✓') else BLACK
    p.font.bold = line.endswith('：') or line.startswith('AI')
    p.font.name = '微软雅黑'

add_rounded(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5), fill_color=WHITE, line_color=GREEN, line_width=Pt(1))
add_text(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5), '⚙️ AI诊断工作流程', size=16, bold=True, color=GREEN)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = ['1️⃣  训练阶段（离线）', '用大量标注好的图像训练模型', '让 AI 学会识别正常 vs 异常', '', '2️⃣  推理阶段（在线）', '输入新患者的超声图像', 'AI 自动分析并输出结果', '', '3️⃣  医生复核', 'AI 结果供医生参考', '最终诊断由医生确认']
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.bold = line.startswith('1') or line.startswith('2') or line.startswith('3')
    p.font.name = '微软雅黑'
print("✅ Slide 16: AI技术原理")


# ============================================================
# Slide 17: 项目荣誉与奖项（新增）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_header(slide, '07 项目荣誉与奖项', '作为项目组长率队获得', 17, TOTAL)

awards = [
    {
        'icon': '🏆',
        'title': '全国AI创新大赛全国总决赛',
        'award': '一等奖',
        'color': GOLD,
        'bg': RGBColor(0xff, 0xf8, 0xe1),
        'border': RGBColor(0xff, 0xb3, 0x00),
        'desc': '全国总决赛最高荣誉',
    },
    {
        'icon': '🥈',
        'title': '全国AI创新大赛医疗应用场景',
        'award': '二等奖',
        'color': SILVER,
        'bg': RGBColor(0xef, 0xf5, 0xf7),
        'border': RGBColor(0x90, 0xa4, 0xae),
        'desc': '医疗赛道优秀表现',
    },
    {
        'icon': '🏅',
        'title': '中国创新方法大赛浙江赛区',
        'award': '优胜奖',
        'color': BRONZE,
        'bg': RGBColor(0xff, 0xf3, 0xe0),
        'border': RGBColor(0xff, 0x8a, 0x65),
        'desc': '浙江省创新方法实践',
    },
]

for i, award in enumerate(awards):
    x = Inches(0.5) + Inches(i * 4.2)
    y = Inches(1.2)
    
    add_rounded(slide, x, y, Inches(3.8), Inches(5.5), fill_color=award['bg'], line_color=award['border'], line_width=Pt(2))
    add_text(slide, x, Inches(1.5), Inches(3.8), Inches(0.8), award['icon'], size=48, bold=True, color=award['color'], align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(1.0), award['title'], size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(0.8), Inches(3.6), Inches(2.2), Inches(0.03), fill_color=award['color'])
    add_rounded(slide, x + Inches(0.8), Inches(3.8), Inches(2.2), Inches(0.7), fill_color=award['color'])
    add_text(slide, x + Inches(0.8), Inches(3.85), Inches(2.2), Inches(0.6), award['award'], size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(4.7), Inches(3.2), Inches(0.5), award['desc'], size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4), '💡 项目组长：管建明 院长 率队参赛', size=12, color=GRAY, align=PP_ALIGN.CENTER)
print("✅ Slide 17: 项目荣誉与奖项")


# ============================================================
# Slide 18: 结尾（复制原Slide 20）
# ============================================================
slide = prs_new.slides.add_slide(prs_new.slide_layouts[6])
set_bg(slide, 0xf8, 0xf9, 0xfa)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=BLUE)
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1), '感谢聆听', size=42, bold=True, color=BLUE)
add_shape(slide, Inches(5.5), Inches(3.2), Inches(3), Inches(0.04), fill_color=BLUE)
add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8), 'Q & A', size=36, bold=True, color=BLACK)
add_rounded(slide, Inches(4.0), Inches(4.8), Inches(5.5), Inches(1.8), fill_color=WHITE, line_color=GRAY_LIGHT, line_width=Pt(1))
add_text(slide, Inches(4.3), Inches(5.0), Inches(5), Inches(0.4), '📧 联系我们', size=16, bold=True, color=BLUE)
add_text(slide, Inches(4.3), Inches(5.5), Inches(5), Inches(0.4), 'AIMED 充盈视界 · 超声造影人工智能联合实验室（第一医院体检中心+南京大学）', size=12, color=GRAY)
add_text(slide, Inches(4.3), Inches(5.9), Inches(5), Inches(0.4), '官网：www.aius.xin    地址：浙江省湖州市', size=12, color=GRAY)
add_shape(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), fill_color=BLUE)
print("✅ Slide 18: 结尾")


# ============================================================
# 保存
# ============================================================
prs_new.save(OUTPUT_FILE)
print(f"\n🎉 PPT 已保存: {OUTPUT_FILE}")
print(f"总页数: {TOTAL}")
print(f"文件大小: {os.path.getsize(OUTPUT_FILE) / 1024 / 1024:.1f} MB")
