#!/usr/bin/env python3
"""
构建优化版 PPT - 口服超声充盈造影对胰腺疾病诊断的 AI 模型
整合 AI 架构图、造影对比图、真实超声图像
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import shutil

OUTPUT_DIR = "/root/.openclaw/workspace/attachments"
ASSETS_DIR = os.path.join(OUTPUT_DIR, "ppt_assets")
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "口服超声充盈造影胰腺疾病AI诊断模型_优化版.pptx")

# 配色
COLORS = {
    'blue': RGBColor(0x1a, 0x73, 0xe8),
    'blue_dark': RGBColor(0x0d, 0x47, 0xa1),
    'green': RGBColor(0x34, 0xa8, 0x53),
    'green_dark': RGBColor(0x1b, 0x5e, 0x20),
    'red': RGBColor(0xea, 0x43, 0x35),
    'orange': RGBColor(0xff, 0x6d, 0x00),
    'purple': RGBColor(0x9c, 0x27, 0xb0),
    'teal': RGBColor(0x00, 0x96, 0x88),
    'gray': RGBColor(0x5f, 0x63, 0x68),
    'gray_light': RGBColor(0xe8, 0xea, 0xed),
    'white': RGBColor(0xff, 0xff, 0xff),
    'black': RGBColor(0x20, 0x21, 0x24),
}

# 幻灯片尺寸 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def set_slide_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=COLORS['black'], alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=12, color=COLORS['gray'], font_name='微软雅黑', line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = line_color if line_color else color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox

def add_header_bar(slide, title, subtitle, slide_num, total_slides):
    # 顶部栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), fill_color=COLORS['blue'])
    add_text_box(slide, Inches(0.8), Inches(0.05), Inches(10), Inches(0.5), title, 
                 font_size=22, bold=True, color=COLORS['white'])
    
    # 副标题
    add_text_box(slide, Inches(0.8), Inches(0.65), Inches(10), Inches(0.35), subtitle,
                 font_size=13, color=COLORS['gray'])
    
    # 页脚
    add_shape(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), fill_color=RGBColor(0xf8, 0xf9, 0xfa))
    add_text_box(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.3),
                 'AIMED 充盈视界 · 超声造影人工智能联合实验室', font_size=9, color=COLORS['gray'])
    add_text_box(slide, Inches(12), Inches(7.12), Inches(1), Inches(0.3),
                 f'{slide_num} / {total_slides}', font_size=9, color=COLORS['gray'], alignment=PP_ALIGN.RIGHT)

def add_image_centered(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(image_path, left, top, width, height)


# ============================================================
# 创建 PPT
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

TOTAL_SLIDES = 20

# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)

# 顶部装饰条
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=COLORS['blue'])

# 主标题
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             '口服超声充盈造影', font_size=36, bold=True, color=COLORS['blue'])
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1),
             '对胰腺疾病诊断的 AI 模型', font_size=36, bold=True, color=COLORS['blue'])

# 副标题
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
             '基于《胃癌超声初筛临床应用中国专家共识意见（2025年版）》诊断逻辑延伸',
             font_size=14, color=COLORS['gray'])

# 分隔线
add_shape(slide, Inches(1.5), Inches(4.2), Inches(3), Inches(0.04), fill_color=COLORS['blue'])

# 信息
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.4),
             '2026 医学人工智能前沿学术会议', font_size=16, color=COLORS['black'])
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
             '湖州第一医院 · 体检中心', font_size=14, color=COLORS['gray'])
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
             '汇报人：XXX 主任    时长：20 分钟    AI 深度：基础入门',
             font_size=12, color=COLORS['gray'])

# 底部
add_shape(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), fill_color=COLORS['blue'])

print("✅ Slide 1: 封面")


# ============================================================
# Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '汇报目录', 'Agenda', 2, TOTAL_SLIDES)

agenda_items = [
    ('01', '临床背景与痛点', '胰腺疾病诊断现状与挑战', COLORS['blue']),
    ('02', '口服超声充盈造影', '技术原理与胰腺成像优势', COLORS['green']),
    ('03', 'AI 模型架构设计', '三阶段诊断流程', COLORS['blue']),
    ('04', 'Su-RADS 胰腺分类', 'AI 辅助分级标准', COLORS['orange']),
    ('05', '充盈造影对比展示', '疾病图像与 AI 标注对比', COLORS['teal']),
    ('06', '临床应用案例', '模拟演示与效果展示', COLORS['purple']),
    ('07', '合规与展望', '科研路径与未来规划', COLORS['red']),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(1.2) + Inches(i * 0.8)
    # 序号
    add_rounded_shape(slide, Inches(1.2), y, Inches(0.6), Inches(0.6), fill_color=color)
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(0.6), Inches(0.45),
                 num, font_size=18, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, Inches(2.0), y + Inches(0.02), Inches(5), Inches(0.35),
                 title, font_size=16, bold=True, color=COLORS['black'])
    # 描述
    add_text_box(slide, Inches(2.0), y + Inches(0.38), Inches(8), Inches(0.3),
                 desc, font_size=11, color=COLORS['gray'])

print("✅ Slide 2: 目录")


# ============================================================
# Slide 3: 临床背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '01 临床背景', '胰腺疾病诊断现状与挑战', 3, TOTAL_SLIDES)

# 左侧 - 流行病学
add_rounded_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
                  fill_color=RGBColor(0xe3, 0xf2, 0xfd), line_color=COLORS['blue'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             '📊 胰腺疾病流行病学', font_size=16, bold=True, color=COLORS['blue'])
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), [
    ('中国胰腺癌发病率逐年上升', False, COLORS['black']),
    ('5 年生存率仅约 10%', False, COLORS['black']),
    ('早期诊断率不足 20%', False, COLORS['black']),
    ('体检人群中胰腺异常检出率约 3-5%', False, COLORS['black']),
    ('中老年人群是高危群体', False, COLORS['black']),
], font_size=13)

# 右侧 - 痛点
add_rounded_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
                  fill_color=RGBColor(0xfc, 0xe4, 0xe4), line_color=COLORS['red'], line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5),
             '⚠️ 当前诊断痛点', font_size=16, bold=True, color=COLORS['red'])
add_multiline_text(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5), [
    ('🔸 胰腺位置深，常规超声成像质量受限', False, COLORS['black']),
    ('🔸 早期病变难以发现，容易漏诊', False, COLORS['black']),
    ('🔸 依赖 CT/MRI，成本高、有辐射', False, COLORS['black']),
    ('🔸 基层医院缺乏经验丰富的超声医师', False, COLORS['black']),
    ('🔸 诊断结果高度依赖操作者经验', False, COLORS['black']),
    ('🔸 缺乏标准化、可量化评估体系', False, COLORS['black']),
], font_size=13)

print("✅ Slide 3: 临床背景")


# ============================================================
# Slide 4: 充盈造影技术
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '02 口服超声充盈造影技术', '技术原理与胰腺成像优势', 4, TOTAL_SLIDES)

# 左侧 - 原理
add_rounded_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(4.5),
                  fill_color=COLORS['white'], line_color=COLORS['blue'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             '🔬 技术原理', font_size=16, bold=True, color=COLORS['blue'])
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(3.5), [
    ('患者口服专用超声造影剂', False, COLORS['black']),
    ('造影剂在胃肠道内形成均匀充盈', False, COLORS['black']),
    ('消除胃肠道气体干扰', False, COLORS['black']),
    ('胃腔作为「声学窗」', False, COLORS['black']),
    ('透过胃窗观察胰腺等邻近器官', False, COLORS['black']),
    ('获得清晰、稳定的胰腺图像', False, COLORS['black']),
], font_size=13)

# 右侧 - 优势
add_rounded_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(4.5),
                  fill_color=COLORS['white'], line_color=COLORS['green'], line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5),
             '✅ 核心优势', font_size=16, bold=True, color=COLORS['green'])
add_multiline_text(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(3.5), [
    ('🟢 无创无辐射，患者接受度高', False, COLORS['black']),
    ('🟢 操作简单，适合体检筛查', False, COLORS['black']),
    ('🟢 成本低廉，可重复检查', False, COLORS['black']),
    ('🟢 成像质量稳定，减少气体干扰', False, COLORS['black']),
    ('🟢 可动态观察胰腺血流情况', False, COLORS['black']),
    ('🟢 适合基层医院推广', False, COLORS['black']),
], font_size=13)

# 底部依据
add_rounded_shape(slide, Inches(0.5), Inches(6.0), Inches(12.5), Inches(0.8),
                  fill_color=RGBColor(0xe8, 0xf5, 0xe9), line_color=COLORS['green'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(6.15), Inches(12), Inches(0.5),
             '💡 技术依据：《胃癌超声初筛临床应用中国专家共识意见（2025年版）》已明确口服超声充盈造影在胃部疾病筛查中的价值和标准化流程',
             font_size=11, color=COLORS['green_dark'])

print("✅ Slide 4: 充盈造影技术")


# ============================================================
# Slide 5: AI 三阶段架构（插入生成的图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '03 AI 模型架构设计', '三阶段智能诊断流程', 5, TOTAL_SLIDES)

# 插入 AI 三阶段流程图
stage_img = os.path.join(ASSETS_DIR, 'ai_3stage_diagram.png')
slide.shapes.add_picture(stage_img, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8))

print("✅ Slide 5: AI 三阶段架构")


# ============================================================
# Slide 6: AI 系统架构全景（插入生成的图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '03 AI 模型架构设计', '系统架构全景图', 6, TOTAL_SLIDES)

# 插入系统架构图
arch_img = os.path.join(ASSETS_DIR, 'ai_system_architecture.png')
slide.shapes.add_picture(arch_img, Inches(0.2), Inches(1.0), Inches(12.9), Inches(5.9))

print("✅ Slide 6: AI 系统架构全景")


# ============================================================
# Slide 7: Su-RADS 分级标准（插入生成的图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '04 Su-RADS 胰腺分类标准', 'AI 辅助分级体系', 7, TOTAL_SLIDES)

# 插入 Su-RADS 图
surads_img = os.path.join(ASSETS_DIR, 'surads_chart.png')
slide.shapes.add_picture(surads_img, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))

print("✅ Slide 7: Su-RADS 分级标准")


# ============================================================
# Slide 8: 充盈造影对比图（插入生成的图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '05 充盈造影对比展示', '疾病充盈造影前后图像对比 + AI 标注', 8, TOTAL_SLIDES)

# 插入对比图
contrast_img = os.path.join(ASSETS_DIR, 'contrast_comparison.png')
slide.shapes.add_picture(contrast_img, Inches(0.2), Inches(1.0), Inches(12.9), Inches(5.8))

print("✅ Slide 8: 充盈造影对比图")


# ============================================================
# Slide 9: 真实超声图像展示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '05 充盈造影对比展示', '真实病例超声图像（充盈后）', 9, TOTAL_SLIDES)

# 选取6张真实超声图像展示
real_images = [
    ('data/cases/stomach/2026-04-23/胃底部恶性间质瘤_A_g_20260423011923.jpg', '胃底部恶性间质瘤'),
    ('data/cases/stomach/2026-04-23/胃角溃疡（活动期）A_o_20260423012114.jpg', '胃角溃疡（活动期）'),
    ('data/cases/stomach/2026-04-23/胃小弯粘膜下良性间质瘤_L_20260423010552.jpg', '胃小弯良性间质瘤'),
    ('data/cases/stomach/2026-04-23/胃体部早期胃 MT 凹陷型 1_1776906580.jpg', '胃体部早期 MT 凹陷型'),
    ('data/cases/stomach/2026-04-23/胃恶性淋巴瘤_aq_20260423012731.jpg', '胃恶性淋巴瘤'),
    ('data/cases/pancreas/2026-04-23/case_20260423000323_62e6efa2.jpg', '胰腺病例'),
]

for i, (img_path, label) in enumerate(real_images):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.2) + Inches(row * 2.8)
    
    full_path = os.path.join('/root/.openclaw/workspace', img_path)
    if os.path.exists(full_path):
        slide.shapes.add_picture(full_path, x, y, Inches(3.8), Inches(2.4))
    
    # 标签
    add_rounded_shape(slide, x, y + Inches(2.4), Inches(3.8), Inches(0.35),
                      fill_color=COLORS['blue'])
    add_text_box(slide, x, y + Inches(2.42), Inches(3.8), Inches(0.3),
                 label, font_size=10, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)

print("✅ Slide 9: 真实超声图像展示")


# ============================================================
# Slide 10: 真实病例对比（充盈前后对比）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '05 充盈造影对比展示', '真实病例对比 - 充盈前后', 10, TOTAL_SLIDES)

# 选取3组对比图像
comparison_groups = [
    [
        ('data/cases/stomach/2026-04-23/胃底部恶性间质瘤（粘膜下型）_c_20260423010519.jpg', '充盈前'),
        ('data/cases/stomach/2026-04-23/胃底部恶性间质瘤_A_g_20260423011923.jpg', '充盈后'),
    ],
    [
        ('data/cases/stomach/2026-04-23/胃角溃疡_R_r_20260423012154.jpg', '充盈前'),
        ('data/cases/stomach/2026-04-23/胃角溃疡（活动期）A_o_20260423012114.jpg', '充盈后'),
    ],
    [
        ('data/cases/stomach/2026-04-23/胃小弯粘膜下良性间质瘤_i_20260423012029.jpg', '充盈前'),
        ('data/cases/stomach/2026-04-23/胃小弯粘膜下良性间质瘤_L_20260423010552.jpg', '充盈后'),
    ],
]

for i, group in enumerate(comparison_groups):
    x = Inches(0.5) + Inches(i * 4.2)
    y = Inches(1.3)
    
    for j, (img_path, label) in enumerate(group):
        full_path = os.path.join('/root/.openclaw/workspace', img_path)
        if os.path.exists(full_path):
            slide.shapes.add_picture(full_path, x, y + Inches(j * 2.5), Inches(3.8), Inches(2.2))
        
        # 标签
        label_color = COLORS['gray'] if label == '充盈前' else COLORS['green']
        add_rounded_shape(slide, x, y + Inches(2.2) + Inches(j * 2.5), Inches(3.8), Inches(0.3),
                          fill_color=label_color)
        add_text_box(slide, x, y + Inches(2.22) + Inches(j * 2.5), Inches(3.8), Inches(0.28),
                     label, font_size=10, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)

print("✅ Slide 10: 真实病例对比")


# ============================================================
# Slide 11: 文档真实对比图（充盈超声造影 vs 胃镜/常规超声）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '05 充盈造影对比展示', '充盈超声造影 vs 胃镜/常规超声 — 真实对比', 11, TOTAL_SLIDES)

# 4张真实对比图 2x2 布局
doc_images = [
    ('/tmp/doc_images/image_0.jpeg', '充盈超声造影 — 胃壁全层清晰显示'),
    ('/tmp/doc_images/image_3.jpeg', '胃镜 — 仅显示黏膜表面'),
    ('/tmp/doc_images/image_12.jpeg', '充盈超声 — 胰腺结构完整显示'),
    ('/tmp/doc_images/image_10.jpeg', '常规超声 — 气体干扰胰腺显示不清'),
]

for i, (img_path, label) in enumerate(doc_images):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.2) + Inches(row * 2.9)
    
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, Inches(5.8), Inches(2.5))
    
    # 标签
    label_color = COLORS['green'] if '充盈' in label else COLORS['gray']
    add_rounded_shape(slide, x, y + Inches(2.5), Inches(5.8), Inches(0.35),
                      fill_color=label_color)
    add_text_box(slide, x, y + Inches(2.52), Inches(5.8), Inches(0.3),
                 label, font_size=10, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)

# 底部说明
add_rounded_shape(slide, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.0),
                  fill_color=None, line_color=None)
add_text_box(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4),
             '💡 来源：《充盈超声造影与胃镜/常规超声对比：理想的普筛早筛工具》',
             font_size=10, color=COLORS['gray'])

print("✅ Slide 11: 文档真实对比图")


# ============================================================
# Slide 12: 标准文档真实超声图像（管建明修订稿）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '05 充盈造影对比展示', '标准文档真实超声图像 — 管建明修订稿', 12, TOTAL_SLIDES)

# 选取标准文档中的6张真实超声图像 3x2 布局
standard_images = [
    ('/tmp/doc2_images/img_1.png', '造影前空腹 vs 造影后充盈对比'),
    ('/tmp/doc2_images/img_5.png', '胰头标准切面（充盈造影后）'),
    ('/tmp/doc2_images/img_6.png', '胰体标准切面（充盈造影后）'),
    ('/tmp/doc2_images/img_13.png', '胰尾长轴切面（充盈造影后）'),
    ('/tmp/doc2_images/img_15.png', '门静脉胆总管切面'),
    ('/tmp/doc2_images/img_17.png', '腹主动脉胰体切面'),
]

for i, (img_path, label) in enumerate(standard_images):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.2) + Inches(row * 2.9)
    
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, Inches(3.8), Inches(2.5))
    
    # 标签
    add_rounded_shape(slide, x, y + Inches(2.5), Inches(3.8), Inches(0.35),
                      fill_color=COLORS['blue'])
    add_text_box(slide, x, y + Inches(2.52), Inches(3.8), Inches(0.3),
                 label, font_size=10, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)

# 底部说明
add_text_box(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4),
             '💡 来源：《标准胰腺超声造影诊断检查标准和 AI 数据采集标准》（管建明修订稿，2025年10月）',
             font_size=10, color=COLORS['gray'])

print("✅ Slide 12: 标准文档真实超声图像")


# ============================================================
# Slide 13: 临床应用案例
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '06 临床应用案例演示', '模拟案例展示 AI 辅助诊断流程', 13, TOTAL_SLIDES)

# 案例1
add_rounded_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
                  fill_color=COLORS['white'], line_color=COLORS['blue'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             '📋 案例一：中年男性，体检筛查', font_size=15, bold=True, color=COLORS['blue'])
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), [
    ('患者信息：', True, COLORS['black']),
    ('男，52 岁，无明显症状', False, COLORS['gray']),
    ('常规体检，有吸烟史', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('检查过程：', True, COLORS['black']),
    ('1. 口服造影剂 500ml', False, COLORS['gray']),
    ('2. 等待 15 分钟充盈', False, COLORS['gray']),
    ('3. 超声检查胰腺', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('AI 分析结果：', True, COLORS['blue']),
    ('Su-RADS 1 · 阴性', False, COLORS['green']),
    ('胰腺形态正常，无异常发现', False, COLORS['gray']),
], font_size=12)

# 案例2
add_rounded_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
                  fill_color=COLORS['white'], line_color=COLORS['orange'], line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5),
             '📋 案例二：老年女性，腹痛就诊', font_size=15, bold=True, color=COLORS['orange'])
add_multiline_text(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5), [
    ('患者信息：', True, COLORS['black']),
    ('女，65 岁，上腹痛 2 周', False, COLORS['gray']),
    ('体重下降，食欲减退', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('检查过程：', True, COLORS['black']),
    ('1. 口服造影剂 500ml', False, COLORS['gray']),
    ('2. 等待 15 分钟充盈', False, COLORS['gray']),
    ('3. 超声检查胰腺', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('AI 分析结果：', True, COLORS['orange']),
    ('Su-RADS 4 · 高危险度', False, COLORS['red']),
    ('可疑恶性，建议增强 CT/MRI', False, COLORS['gray']),
], font_size=12)

print("✅ Slide 11: 临床应用案例")


# ============================================================
# Slide 12: 合规边界
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '07 合规边界与科研工具定位', 'Phase 1 科研工具定位', 14, TOTAL_SLIDES)

# 允许
add_rounded_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
                  fill_color=RGBColor(0xe8, 0xf5, 0xe9), line_color=COLORS['green'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             '✅ 当前阶段允许', font_size=16, bold=True, color=COLORS['green'])
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), [
    ('✓ 科研合作与技术验证', False, COLORS['black']),
    ('✓ 学术交流与论文发表', False, COLORS['black']),
    ('✓ 合规数据采集与标注', False, COLORS['black']),
    ('✓ 内部测试与性能评估', False, COLORS['black']),
    ('✓ 与三甲医院合作研究', False, COLORS['black']),
    ('✓ 伦理审查后的临床试验', False, COLORS['black']),
    ('✓ 技术演示与教学培训', False, COLORS['black']),
], font_size=13)

# 禁止
add_rounded_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
                  fill_color=RGBColor(0xfc, 0xe4, 0xe4), line_color=COLORS['red'], line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5),
             '⚠️ 当前阶段禁止', font_size=16, bold=True, color=COLORS['red'])
add_multiline_text(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5), [
    ('✗ 对外宣称「诊断」功能', False, COLORS['black']),
    ('✗ 直接面向患者收费服务', False, COLORS['black']),
    ('✗ 未经伦理审批的数据使用', False, COLORS['black']),
    ('✗ 商业化收费运营', False, COLORS['black']),
    ('✗ 替代医生独立出具报告', False, COLORS['black']),
    ('✗ 未通过三类医疗器械审批', False, COLORS['black']),
    ('✗ 超出科研范围的应用', False, COLORS['black']),
], font_size=13)

print("✅ Slide 13: 合规边界")


# ============================================================
# Slide 13: 发展路线图（插入生成的图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '07 发展路线图', '三阶段推进策略', 15, TOTAL_SLIDES)

# 插入路线图
roadmap_img = os.path.join(ASSETS_DIR, 'roadmap.png')
slide.shapes.add_picture(roadmap_img, Inches(0.2), Inches(1.0), Inches(12.9), Inches(5.8))

print("✅ Slide 14: 发展路线图")


# ============================================================
# Slide 14: 核心优势（插入生成的图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '核心优势总结', '为什么选择口服超声充盈造影 + AI', 16, TOTAL_SLIDES)

# 插入核心优势图
adv_img = os.path.join(ASSETS_DIR, 'advantages_chart.png')
slide.shapes.add_picture(adv_img, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))

print("✅ Slide 15: 核心优势")


# ============================================================
# Slide 15: 技术架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '技术架构概览', '系统如何工作（非技术视角）', 17, TOTAL_SLIDES)

layers = [
    ('🖥️ 用户界面层', '医生/患者操作界面', COLORS['blue'], Inches(1.3)),
    ('🧠 AI 分析层', '图像识别 + 智能诊断', COLORS['green'], Inches(2.7)),
    ('📊 数据层', '病例数据库 + 知识库', COLORS['orange'], Inches(4.1)),
    ('🔐 存证层', '区块链存证 + 审计', COLORS['purple'], Inches(5.5)),
]

for title, desc, color, y in layers:
    add_rounded_shape(slide, Inches(1.5), y, Inches(10), Inches(1.1),
                      fill_color=COLORS['white'], line_color=color, line_width=Pt(2))
    add_text_box(slide, Inches(2.0), y + Inches(0.15), Inches(4), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(slide, Inches(6.0), y + Inches(0.35), Inches(5), Inches(0.4),
                 desc, font_size=13, color=COLORS['gray'])
    # 箭头
    if color != COLORS['purple']:
        next_y = y + Inches(1.3)
        add_text_box(slide, Inches(11.5), y + Inches(0.3), Inches(0.5), Inches(0.5),
                     '↓', font_size=20, color=color, alignment=PP_ALIGN.CENTER)

# 底部
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             '💡 所有数据本地化部署，符合医疗数据安全要求', font_size=11, color=COLORS['gray'])

print("✅ Slide 16: 技术架构")


# ============================================================
# Slide 16: AI 技术原理
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, 'AI 技术原理（基础入门）', '无需编程知识，理解 AI 如何辅助诊断', 18, TOTAL_SLIDES)

# 左侧
add_rounded_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
                  fill_color=COLORS['white'], line_color=COLORS['blue'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             '🤖 AI 在医学影像中的角色', font_size=16, bold=True, color=COLORS['blue'])
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), [
    ('AI = 计算机视觉 + 深度学习', True, COLORS['black']),
    ('', False, COLORS['gray']),
    ('类比理解：', True, COLORS['black']),
    ('就像培养一个「超级助手」', False, COLORS['gray']),
    ('看过 10000+ 张超声图像后', False, COLORS['gray']),
    ('能自动识别可疑病灶', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('AI 不是替代医生，而是：', True, COLORS['black']),
    ('✓ 提高诊断效率', False, COLORS['green']),
    ('✓ 减少漏诊误诊', False, COLORS['green']),
    ('✓ 统一诊断标准', False, COLORS['green']),
], font_size=13)

# 右侧
add_rounded_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
                  fill_color=COLORS['white'], line_color=COLORS['green'], line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5),
             '⚙️ AI诊断工作流程', font_size=16, bold=True, color=COLORS['green'])
add_multiline_text(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5), [
    ('1️⃣  训练阶段（离线）', True, COLORS['black']),
    ('用大量标注好的图像训练模型', False, COLORS['gray']),
    ('让 AI 学会识别正常 vs 异常', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('2️⃣  推理阶段（在线）', True, COLORS['black']),
    ('输入新患者的超声图像', False, COLORS['gray']),
    ('AI 自动分析并输出结果', False, COLORS['gray']),
    ('', False, COLORS['gray']),
    ('3️⃣  医生复核', True, COLORS['black']),
    ('AI 结果供医生参考', False, COLORS['gray']),
    ('最终诊断由医生确认', False, COLORS['gray']),
], font_size=13)

print("✅ Slide 17: AI 技术原理")


# ============================================================
# Slide 17: 未来展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)
add_header_bar(slide, '未来展望', '从胰腺到多病种的 AI诊断平台', 19, TOTAL_SLIDES)

# 短期
add_rounded_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
                  fill_color=RGBColor(0xe3, 0xf2, 0xfd), line_color=COLORS['blue'], line_width=Pt(1))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             '🔮 短期目标（1-2 年）', font_size=16, bold=True, color=COLORS['blue'])
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), [
    ('✓ 完成胰腺疾病AI 模型验证', False, COLORS['black']),
    ('✓ 积累 1000+ 标注病例', False, COLORS['black']),
    ('✓ 发表学术论文 2-3 篇', False, COLORS['black']),
    ('✓ 建立 3-5 家合作医院', False, COLORS['black']),
    ('✓ 通过伦理审查', False, COLORS['black']),
    ('✓ 完成 Phase 1 技术验证', False, COLORS['black']),
], font_size=13)

# 长期
add_rounded_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5),
                  fill_color=RGBColor(0xf3, 0xe5, 0xf5), line_color=COLORS['purple'], line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.5),
             '🌟 长期愿景（3-5 年）', font_size=16, bold=True, color=COLORS['purple'])
add_multiline_text(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.5), [
    ('🔸 扩展至胃、肝、胆等多病种', False, COLORS['black']),
    ('🔸 建立国家级超声 AI 数据库', False, COLORS['black']),
    ('🔸 推动基层医院普及超声筛查', False, COLORS['black']),
    ('🔸 实现「三甲医院能力下沉」', False, COLORS['black']),
    ('🔸 申报 NMPA 三类医疗器械', False, COLORS['black']),
    ('🔸 探索医保覆盖与商业化路径', False, COLORS['black']),
], font_size=13)

print("✅ Slide 19: 未来展望")


# ============================================================
# Slide 18: 结尾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, 0xf8, 0xf9, 0xfa)

# 顶部装饰条
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=COLORS['blue'])

# 感谢
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
             '感谢聆听', font_size=42, bold=True, color=COLORS['blue'])

# 分隔线
add_shape(slide, Inches(5.5), Inches(3.2), Inches(3), Inches(0.04), fill_color=COLORS['blue'])

# Q&A
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
             'Q & A', font_size=36, bold=True, color=COLORS['black'])

# 联系方式
add_rounded_shape(slide, Inches(4.0), Inches(4.8), Inches(5.5), Inches(1.8),
                  fill_color=COLORS['white'], line_color=COLORS['gray_light'], line_width=Pt(1))
add_text_box(slide, Inches(4.3), Inches(5.0), Inches(5), Inches(0.4),
             '📧 联系我们', font_size=16, bold=True, color=COLORS['blue'])
add_text_box(slide, Inches(4.3), Inches(5.5), Inches(5), Inches(0.4),
             'AIMED 充盈视界 · 超声造影人工智能联合实验室', font_size=12, color=COLORS['gray'])
add_text_box(slide, Inches(4.3), Inches(5.9), Inches(5), Inches(0.4),
             '官网：www.aius.xin    地址：浙江省湖州市', font_size=12, color=COLORS['gray'])

# 底部
add_shape(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), fill_color=COLORS['blue'])

print("✅ Slide 20: 结尾")


# ============================================================
# 保存
# ============================================================
prs.save(OUTPUT_FILE)
print(f"\n🎉 PPT 已保存: {OUTPUT_FILE}")
print(f"文件大小: {os.path.getsize(OUTPUT_FILE) / 1024 / 1024:.1f} MB")
