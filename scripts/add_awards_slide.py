#!/usr/bin/env python3
"""
在 PPT 末尾（感谢页之前）插入奖项页
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation('/tmp/管院长-正式稿.pptx')
total = len(prs.slides)

# 配色
BLUE = RGBColor(0x1a, 0x73, 0xe8)
BLUE_DARK = RGBColor(0x0d, 0x47, 0xa1)
GOLD = RGBColor(0xff, 0xb3, 0x00)
GOLD_DARK = RGBColor(0xcc, 0x88, 0x00)
SILVER = RGBColor(0x90, 0xa4, 0xae)
BRONZE = RGBColor(0xff, 0x8a, 0x65)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x20, 0x21, 0x24)
GRAY = RGBColor(0x5f, 0x63, 0x68)
GRAY_LIGHT = RGBColor(0xe8, 0xea, 0xed)

def set_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return txBox

def add_header(slide, title, subtitle, num, total):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.6), fill_color=BLUE)
    add_text(slide, Inches(0.8), Inches(0.05), Inches(10), Inches(0.5), title, size=22, bold=True, color=WHITE)
    add_text(slide, Inches(0.8), Inches(0.65), Inches(10), Inches(0.35), subtitle, size=13, color=GRAY)
    add_shape(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), fill_color=RGBColor(0xf8, 0xf9, 0xfa))
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.3), 'AIMED 充盈视界 · 超声造影人工智能联合实验室（第一医院体检中心+南京大学）', size=9, color=GRAY)
    add_text(slide, Inches(12), Inches(7.12), Inches(1), Inches(0.3), f'{num} / {total}', size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# 插入新幻灯片（在最后一页之前）
new_slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(new_slide, 0xf8, 0xf9, 0xfa)
new_slide_num = total  # 插入在倒数第二页（感谢页之前）
new_total = total + 1

add_header(new_slide, '项目荣誉与奖项', '作为项目组长率队获得', new_slide_num, new_total)

# 奖项数据
awards = [
    {
        'icon': '🏆',
        'title': '全国AI创新大赛全国总决赛',
        'award': '一等奖',
        'color': GOLD,
        'bg': RGBColor(0xff, 0xf8, 0xe1),
        'border': RGBColor(0xff, 0xb3, 0x00),
        'desc': '全国总决赛最高荣誉',
    },
    {
        'icon': '🥈',
        'title': '全国AI创新大赛医疗应用场景',
        'award': '二等奖',
        'color': SILVER,
        'bg': RGBColor(0xef, 0xf5, 0xf7),
        'border': RGBColor(0x90, 0xa4, 0xae),
        'desc': '医疗赛道优秀表现',
    },
    {
        'icon': '🏅',
        'title': '中国创新方法大赛浙江赛区',
        'award': '优胜奖',
        'color': BRONZE,
        'bg': RGBColor(0xff, 0xf3, 0xe0),
        'border': RGBColor(0xff, 0x8a, 0x65),
        'desc': '浙江省创新方法实践',
    },
]

for i, award in enumerate(awards):
    x = Inches(0.5) + Inches(i * 4.2)
    y = Inches(1.2)
    
    # 卡片
    add_rounded(new_slide, x, y, Inches(3.8), Inches(5.5), fill_color=award['bg'], line_color=award['border'], line_width=Pt(2))
    
    # 图标
    add_text(new_slide, x, Inches(1.5), Inches(3.8), Inches(0.8), award['icon'], size=48, bold=True, color=award['color'], align=PP_ALIGN.CENTER)
    
    # 奖项名称
    add_text(new_slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(1.0), award['title'], size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    
    # 分隔线
    add_shape(new_slide, x + Inches(0.8), Inches(3.6), Inches(2.2), Inches(0.03), fill_color=award['color'])
    
    # 奖项等级
    add_rounded(new_slide, x + Inches(0.8), Inches(3.8), Inches(2.2), Inches(0.7), fill_color=award['color'])
    add_text(new_slide, x + Inches(0.8), Inches(3.85), Inches(2.2), Inches(0.6), award['award'], size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    # 描述
    add_text(new_slide, x + Inches(0.3), Inches(4.7), Inches(3.2), Inches(0.5), award['desc'], size=12, color=GRAY, align=PP_ALIGN.CENTER)

# 底部说明
add_text(new_slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4),
         '💡 项目组长：管建明 院长 率队参赛', size=12, color=GRAY, align=PP_ALIGN.CENTER)

# 更新所有后续幻灯片的页码
for i in range(new_slide_num, total):
    slide = prs.slides[i]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text
                # 更新页码格式 "X / 20" -> "X / 21"
                if '/ 20' in text:
                    parts = text.split('/ 20')
                    num = int(parts[0].strip())
                    if num >= new_slide_num:
                        num += 1
                    para.text = f'{num} / {new_total}'

# 更新目录页
for i, slide in enumerate(prs.slides):
    if i == 1:  # 目录页
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if '/ 20' in para.text:
                        para.text = para.text.replace('/ 20', f'/ {new_total}')

# 保存
output_path = '/root/.openclaw/workspace/attachments/管院长-正式稿_含奖项.pptx'
prs.save(output_path)
print(f"✅ 奖项页已添加，共 {new_total} 页")
print(f"文件: {output_path}")
print(f"大小: {os.path.getsize(output_path) / 1024 / 1024:.1f} MB")
