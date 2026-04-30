#!/usr/bin/env python3
"""
生成 AIMED 口服超声充盈造影胰腺疾病 AI 诊断模型 PPT
用途：医学行业 AI 学术会议
时长：20 分钟
AI 深度：基础
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================
# 配色方案 - 医疗专业风格
# ============================================
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)       # 深蓝
SECONDARY = RGBColor(0x02, 0x84, 0xC7)     # 亮蓝
ACCENT = RGBColor(0x05, 0x96, 0x69)        # 绿色
DARK = RGBColor(0x1E, 0x29, 0x3B)          # 深灰
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)         # 浅蓝灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xDC, 0x26, 0x26)
GRAY = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================
# 辅助函数
# ============================================

def add_background(slide, color):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """添加形状并返回"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = font_name
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, font_color=DARK, spacing=Pt(8)):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.space_after = spacing
        p.level = 0
        p.font.name = "微软雅黑"
    return txBox

def add_header_bar(slide, title, subtitle=""):
    """添加顶部标题栏"""
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.7), title, font_size=32, font_color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.4), subtitle, font_size=16, font_color=RGBColor(0xBF, 0xDB, 0xFE))

def add_footer(slide, page_num, total=20):
    """添加页脚"""
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.4), 
                 "AIMED 充盈视界 · 超声造影人工智能联合实验室", font_size=10, font_color=GRAY)
    add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                 f"{page_num} / {total}", font_size=10, font_color=GRAY, alignment=PP_ALIGN.RIGHT)

# ============================================
# Slide 1: 封面
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# 顶部装饰条
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), PRIMARY)

# 主标题
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5),
             "口服超声充盈造影对胰腺疾病诊断的\nAI 模型初构",
             font_size=40, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
             "基于《胃癌超声初筛临床应用中国专家共识意见（2025年版）》诊断逻辑延伸",
             font_size=20, font_color=GRAY, alignment=PP_ALIGN.CENTER)

# 会议信息
add_shape(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(1.2), LIGHT)
add_text_box(slide, Inches(3.5), Inches(4.9), Inches(6.3), Inches(0.5),
             "2026 医学人工智能前沿学术会议", font_size=18, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(5.4), Inches(6.3), Inches(0.5),
             "湖州第一医院 · 体检中心", font_size=16, font_color=DARK, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.5),
             "汇报人：XXX 主任    时长：20 分钟    AI 深度：基础入门",
             font_size=14, font_color=GRAY, alignment=PP_ALIGN.CENTER)

# 底部装饰条
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), PRIMARY)

# ============================================
# Slide 2: 目录
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "汇报目录", "Agenda")
add_footer(slide, 2)

toc_items = [
    ("01", "临床背景与痛点", "胰腺疾病诊断现状与挑战"),
    ("02", "口服超声充盈造影", "技术原理与胰腺成像优势"),
    ("03", "AI 模型架构设计", "三阶段诊断流程"),
    ("04", "Su-RADS 胰腺分类", "AI 辅助分级标准"),
    ("05", "临床应用案例", "模拟演示与效果展示"),
    ("06", "合规与展望", "科研路径与未来规划"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(i * 0.9)
    
    # 序号圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY if i < 3 else SECONDARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(2.1), y, Inches(4), Inches(0.35), title, font_size=20, font_color=DARK, bold=True)
    add_text_box(slide, Inches(2.1), y + Inches(0.35), Inches(8), Inches(0.3), desc, font_size=13, font_color=GRAY)

# ============================================
# Slide 3: 临床背景 - 胰腺疾病现状
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "01 临床背景", "胰腺疾病诊断现状与挑战")
add_footer(slide, 3)

# 左侧：统计数据
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.7), Inches(0.5),
             "📊 胰腺疾病流行病学", font_size=22, font_color=PRIMARY, bold=True)

stats = [
    "中国胰腺癌发病率逐年上升",
    "5 年生存率仅约 10%",
    "早期诊断率不足 20%",
    "体检人群中胰腺异常检出率约 3-5%",
    "中老年人群是高危群体",
]
add_bullet_list(slide, Inches(1.2), Inches(2.5), Inches(4.7), Inches(4), stats, font_size=15, spacing=Pt(12))

# 右侧：痛点
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
             "⚠️ 当前诊断痛点", font_size=22, font_color=RED, bold=True)

pain_points = [
    "🔸 胰腺位置深，常规超声成像质量受限",
    "🔸 早期病变难以发现，容易漏诊",
    "🔸 依赖 CT/MRI，成本高、有辐射",
    "🔸 基层医院缺乏经验丰富的超声医师",
    "🔸 诊断结果高度依赖操作者经验",
    "🔸 缺乏标准化、可量化的诊断体系",
]
add_bullet_list(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(4), pain_points, font_size=15, spacing=Pt(12))

# ============================================
# Slide 4: 口服超声充盈造影
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "02 口服超声充盈造影技术", "技术原理与胰腺成像优势")
add_footer(slide, 4)

# 左侧：技术原理
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "🔬 技术原理", font_size=22, font_color=PRIMARY, bold=True)

principle = [
    "患者口服专用超声造影剂",
    "造影剂在胃肠道内形成均匀充盈",
    "消除胃肠道气体干扰",
    "胃腔作为「声学窗」",
    "透过胃窗观察胰腺等邻近器官",
    "获得清晰、稳定的胰腺图像",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5), principle, font_size=15, spacing=Pt(10))

# 右侧：优势对比
add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.5),
             "✅ 核心优势", font_size=22, font_color=ACCENT, bold=True)

advantages = [
    "🟢 无创无辐射，患者接受度高",
    "🟢 操作简单，适合体检筛查",
    "🟢 成本低廉，可重复检查",
    "🟢 成像质量稳定，减少气体干扰",
    "🟢 可动态观察胰腺血流情况",
    "🟢 适合基层医院推广",
]
add_bullet_list(slide, Inches(6.8), Inches(2.3), Inches(5.8), Inches(3.5), advantages, font_size=15, spacing=Pt(10))

# 底部：共识依据
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.2), PRIMARY)
add_text_box(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.9),
             "💡 技术依据：《胃癌超声初筛临床应用中国专家共识意见（2025年版）》\n已明确口服超声充盈造影在胃部疾病筛查中的价值和标准化流程",
             font_size=14, font_color=WHITE)

# ============================================
# Slide 5: AI 模型架构 - 总览
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "03 AI 模型架构设计", "三阶段智能诊断流程")
add_footer(slide, 5)

# 架构图
# 阶段 1
add_shape(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(4.5), RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(2.7), Inches(0.5),
             "阶段一：数据采集与预处理", font_size=18, font_color=PRIMARY, bold=True)

stage1 = [
    "📷 图像/视频采集",
    "   - 超声图像质量评估",
    "   - 自动去噪增强",
    "   - 图像标准化处理",
    "",
    "📋 患者信息整合",
    "   - 年龄、病史、风险因素",
    "   - 高危人群自动识别",
]
add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(2.7), Inches(3.5), stage1, font_size=13, spacing=Pt(6))

# 箭头
arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(3.8), Inches(0.5), Inches(0.3))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = PRIMARY
arrow1.line.fill.background()

# 阶段 2
add_shape(slide, Inches(5.2), Inches(2.0), Inches(3.5), Inches(4.5), RGBColor(0xF0, 0xFD, 0xF4))
add_text_box(slide, Inches(5.6), Inches(2.2), Inches(2.7), Inches(0.5),
             "阶段二：AI 智能分析", font_size=18, font_color=ACCENT, bold=True)

stage2 = [
    "🔍 病灶自动检测",
    "   - 胰腺轮廓自动勾画",
    "   - 可疑病灶区域标注",
    "",
    "📐 特征量化测量",
    "   - 胰腺大小、形态",
    "   - 胰管直径测量",
    "   - 回声特征分析",
    "",
    "🧠 智能分级评估",
    "   - Su-RADS 胰腺分类",
    "   - 恶性风险概率计算",
]
add_bullet_list(slide, Inches(5.6), Inches(2.8), Inches(2.7), Inches(3.5), stage2, font_size=13, spacing=Pt(6))

# 箭头
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.9), Inches(3.8), Inches(0.5), Inches(0.3))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = PRIMARY
arrow2.line.fill.background()

# 阶段 3
add_shape(slide, Inches(9.6), Inches(2.0), Inches(3.0), Inches(4.5), RGBColor(0xFF, 0xFB, 0xEB))
add_text_box(slide, Inches(10.0), Inches(2.2), Inches(2.2), Inches(0.5),
             "阶段三：报告生成", font_size=18, font_color=RGBColor(0xD9, 0x77, 0x06), bold=True)

stage3 = [
    "📄 结构化报告",
    "   - 自动填充诊断结论",
    "   - 分级建议处理方案",
    "",
    "👨‍⚕️ 医生复核",
    "   - AI 标注人工确认",
    "   - 补充临床意见",
    "",
    "🔐 区块链存证",
    "   - 报告不可篡改",
    "   - 全程可追溯",
]
add_bullet_list(slide, Inches(10.0), Inches(2.8), Inches(2.2), Inches(3.5), stage3, font_size=13, spacing=Pt(6))

# ============================================
# Slide 6: Su-RADS 胰腺分类
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "04 Su-RADS 胰腺分类标准", "AI 辅助分级体系")
add_footer(slide, 6)

# 分类表格
surads = [
    ("Su-RADS 1", "阴性", "胰腺形态正常，无异常发现", RGBColor(0x05, 0x96, 0x69)),
    ("Su-RADS 2", "低危险度", "良性病变，恶性风险 <2%", RGBColor(0x02, 0x84, 0xC7)),
    ("Su-RADS 3", "中危险度", "可能良性，建议短期复查", RGBColor(0xD9, 0x77, 0x06)),
    ("Su-RADS 4", "高危险度", "可疑恶性，建议进一步检查", RGBColor(0xEA, 0x58, 0x0C)),
    ("Su-RADS 5", "极高危险度", "高度怀疑恶性，建议活检", RGBColor(0xDC, 0x26, 0x26)),
]

for i, (level, desc, detail, color) in enumerate(surads):
    y = Inches(1.6) + Inches(i * 1.0)
    
    # 等级标签
    add_shape(slide, Inches(0.8), y, Inches(2.0), Inches(0.85), color)
    add_text_box(slide, Inches(0.8), y + Inches(0.15), Inches(2.0), Inches(0.55),
                 level, font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_text_box(slide, Inches(3.0), y + Inches(0.1), Inches(3.0), Inches(0.35),
                 desc, font_size=16, font_color=DARK, bold=True)
    
    # 详细说明
    add_text_box(slide, Inches(3.0), y + Inches(0.45), Inches(9.5), Inches(0.35),
                 detail, font_size=13, font_color=GRAY)

# 底部说明
add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.6), LIGHT)
add_text_box(slide, Inches(1.2), Inches(6.65), Inches(11), Inches(0.5),
             "💡 参考 BI-RADS 乳腺分类体系，结合胰腺超声造影特征制定",
             font_size=13, font_color=DARK)

# ============================================
# Slide 7: AI 技术原理（基础版）
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "AI 技术原理（基础入门）", "无需编程知识，理解 AI 如何辅助诊断")
add_footer(slide, 7)

# 左侧：AI 是什么
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "🤖 AI 在医学影像中的角色", font_size=22, font_color=PRIMARY, bold=True)

ai_intro = [
    "AI = 计算机视觉 + 深度学习",
    "",
    "类比理解：",
    "  就像培养一个「超级助手」",
    "  看过 10000+ 张超声图像后",
    "  能自动识别可疑病灶",
    "",
    "AI 不是替代医生，而是：",
    "  ✓ 提高诊断效率",
    "  ✓ 减少漏诊误诊",
    "  ✓ 统一诊断标准",
    "  ✓ 辅助基层医生",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), ai_intro, font_size=15, spacing=Pt(8))

# 右侧：工作流程
add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.5),
             "⚙️ AI 诊断工作流程", font_size=22, font_color=SECONDARY, bold=True)

workflow = [
    "1️⃣  训练阶段（离线）",
    "   用大量标注好的图像训练模型",
    "   让 AI 学会识别正常 vs 异常",
    "",
    "2️⃣  推理阶段（在线）",
    "   输入新患者的超声图像",
    "   AI 自动分析并输出结果",
    "",
    "3️⃣  医生复核（关键环节）",
    "   医生确认 AI 的标注和建议",
    "   结合临床经验做出最终判断",
    "",
    "⚠️  重要提示：",
    "   AI 辅助诊断 ≠ AI 独立诊断",
    "   最终决策权始终在医生手中",
]
add_bullet_list(slide, Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5), workflow, font_size=14, spacing=Pt(6))

# ============================================
# Slide 8: 模型训练数据
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "训练数据与模型性能", "数据质量决定 AI 诊断上限")
add_footer(slide, 8)

# 左侧：数据来源
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.7), Inches(0.5),
             "📚 训练数据来源", font_size=22, font_color=PRIMARY, bold=True)

data_sources = [
    "🏥 合作医院超声图像库",
    "   - 胰腺超声造影图像",
    "   - 标注：正常/良性/恶性",
    "",
    "📖 公开医学数据集",
    "   - 胰腺超声标准图像",
    "   - 国际多中心研究数据",
    "",
    "👨‍⚕️ 专家标注团队",
    "   - 三甲医院超声医师",
    "   - 双盲标注，确保质量",
    "",
    "🔒 数据合规",
    "   - 脱敏处理，保护隐私",
    "   - 伦理审查通过",
]
add_bullet_list(slide, Inches(1.2), Inches(2.5), Inches(4.7), Inches(4), data_sources, font_size=14, spacing=Pt(8))

# 右侧：性能指标
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
             "📊 目标性能指标", font_size=22, font_color=SECONDARY, bold=True)

metrics = [
    "灵敏度（Sensitivity）",
    "   目标：≥ 85%",
    "   含义：不漏诊，减少假阴性",
    "",
    "特异度（Specificity）",
    "   目标：≥ 80%",
    "   含义：不误诊，减少假阳性",
    "",
    "诊断一致性（Kappa）",
    "   目标：≥ 0.75",
    "   含义：AI 与专家意见一致",
    "",
    "诊断时间",
    "   目标：≤ 3 分钟/例",
    "   含义：大幅提高效率",
]
add_bullet_list(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(4), metrics, font_size=14, spacing=Pt(8))

# ============================================
# Slide 9: 临床应用案例
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "05 临床应用案例演示", "模拟案例展示 AI 辅助诊断流程")
add_footer(slide, 9)

# 案例卡片
case1_y = Inches(1.6)
add_shape(slide, Inches(0.8), case1_y, Inches(5.8), Inches(5.2), LIGHT)
add_text_box(slide, Inches(1.2), case1_y + Inches(0.2), Inches(5), Inches(0.5),
             "📋 案例一：中年男性，体检筛查", font_size=20, font_color=PRIMARY, bold=True)

case1 = [
    "患者信息：",
    "  男，52 岁，无明显症状",
    "  常规体检，有吸烟史",
    "",
    "检查过程：",
    "  1. 口服造影剂 500ml",
    "  2. 等待 15 分钟充盈",
    "  3. 超声检查胰腺",
    "",
    "AI 分析结果：",
    "  Su-RADS 2 类（低危险度）",
    "  胰腺形态正常，胰管轻度扩张",
    "  建议：6 个月后复查",
    "",
    "医生复核：",
    "  同意 AI 结论，建议定期随访",
]
add_bullet_list(slide, Inches(1.2), case1_y + Inches(0.8), Inches(5), Inches(4.2), case1, font_size=13, spacing=Pt(6))

# 案例 2
case2_y = Inches(1.6)
add_shape(slide, Inches(6.8), case2_y, Inches(5.8), Inches(5.2), RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, Inches(7.2), case2_y + Inches(0.2), Inches(5), Inches(0.5),
             "📋 案例二：老年女性，腹痛就诊", font_size=20, font_color=RED, bold=True)

case2 = [
    "患者信息：",
    "  女，65 岁，上腹痛 2 周",
    "  体重下降，食欲减退",
    "",
    "检查过程：",
    "  1. 口服造影剂 500ml",
    "  2. 等待 15 分钟充盈",
    "  3. 超声检查胰腺",
    "",
    "AI 分析结果：",
    "  Su-RADS 4 类（高危险度）",
    "  胰头部低回声占位，可疑恶性",
    "  建议：增强 CT + 穿刺活检",
    "",
    "医生复核：",
    "  同意 AI 结论，安排进一步检查",
    "  病理确诊：胰腺导管腺癌",
]
add_bullet_list(slide, Inches(7.2), case2_y + Inches(0.8), Inches(5), Inches(4.2), case2, font_size=13, spacing=Pt(6))

# ============================================
# Slide 10: 合规边界
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "06 合规边界与科研路径", " Phase 1 科研工具定位")
add_footer(slide, 10)

# 左侧：允许做的事
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xF0, 0xFD, 0xF4))
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.7), Inches(0.5),
             "✅ 当前阶段允许", font_size=22, font_color=ACCENT, bold=True)

allowed = [
    "✓ 科研合作与技术验证",
    "✓ 学术交流与论文发表",
    "✓ 合规数据采集与标注",
    "✓ 内部测试与性能评估",
    "✓ 与三甲医院合作研究",
    "✓ 伦理审查后的临床试验",
    "✓ 技术演示与教学培训",
]
add_bullet_list(slide, Inches(1.2), Inches(2.5), Inches(4.7), Inches(4), allowed, font_size=15, spacing=Pt(12))

# 右侧：禁止做的事
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
             "⚠️ 当前阶段禁止", font_size=22, font_color=RED, bold=True)

not_allowed = [
    "✗ 对外宣称「诊断」功能",
    "✗ 直接面向患者收费服务",
    "✗ 未经伦理审批的数据使用",
    "✗ 商业化收费运营",
    "✗ 替代医生独立出具报告",
    "✗ 未通过三类医疗器械审批",
    "✗ 超出科研范围的应用",
]
add_bullet_list(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(4), not_allowed, font_size=15, spacing=Pt(12))

# ============================================
# Slide 11: 三阶段路线图
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "发展路线图", "三阶段推进策略")
add_footer(slide, 11)

# 阶段 1
add_shape(slide, Inches(0.8), Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(3), Inches(0.5),
             "Phase 1（当前）", font_size=22, font_color=PRIMARY, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(3), Inches(0.4),
             "2026 Q2-Q4", font_size=14, font_color=GRAY)

p1_items = [
    "🎯 定位：科研工具",
    "",
    "✓ 技术验证与优化",
    "✓ 积累 500+ 标注病例",
    "✓ AI 准确率 ≥ 85%",
    "✓ 科研合作模式",
    "✓ 伦理审查通过",
    "✓ 内部测试验证",
]
add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(3), Inches(3.5), p1_items, font_size=14, spacing=Pt(8))

# 阶段 2
add_shape(slide, Inches(4.9), Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0xF0, 0xFD, 0xF4))
add_text_box(slide, Inches(5.3), Inches(1.8), Inches(3), Inches(0.5),
             "Phase 2", font_size=22, font_color=ACCENT, bold=True)
add_text_box(slide, Inches(5.3), Inches(2.3), Inches(3), Inches(0.4),
             "2027 Q1-Q3", font_size=14, font_color=GRAY)

p2_items = [
    "🎯 定位：产品化",
    "",
    "✓ 二类器械备案",
    "✓ 试点医院验证",
    "✓ 标准化操作流程",
    "✓ 多中心临床研究",
    "✓ 医生培训体系",
    "✓ 商业模式探索",
]
add_bullet_list(slide, Inches(5.3), Inches(2.8), Inches(3), Inches(3.5), p2_items, font_size=14, spacing=Pt(8))

# 阶段 3
add_shape(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0xFF, 0xFB, 0xEB))
add_text_box(slide, Inches(9.4), Inches(1.8), Inches(3), Inches(0.5),
             "Phase 3", font_size=22, font_color=RGBColor(0xD9, 0x77, 0x06), bold=True)
add_text_box(slide, Inches(9.4), Inches(2.3), Inches(3), Inches(0.4),
             "2027 Q4+", font_size=14, font_color=GRAY)

p3_items = [
    "🎯 定位：规模化",
    "",
    "✓ NMPA 三类证申报",
    "✓ 规模化临床应用",
    "✓ 医保覆盖探索",
    "✓ 基层医院推广",
    "✓ 多病种扩展",
    "✓ 国际化布局",
]
add_bullet_list(slide, Inches(9.4), Inches(2.8), Inches(3), Inches(3.5), p3_items, font_size=14, spacing=Pt(8))

# ============================================
# Slide 12: 技术架构（简化版）
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "技术架构概览", "系统如何工作（非技术视角）")
add_footer(slide, 12)

# 架构分层
layers = [
    ("🖥️ 用户界面层", "医生/患者操作界面", RGBColor(0xEF, 0xF6, 0xFF)),
    ("🧠 AI 分析层", "图像识别 + 智能诊断", RGBColor(0xF0, 0xFD, 0xF4)),
    ("📊 数据层", "病例数据库 + 知识库", RGBColor(0xFF, 0xFB, 0xEB)),
    ("🔐 存证层", "区块链存证 + 审计", RGBColor(0xFE, 0xF2, 0xF2)),
]

for i, (layer, desc, color) in enumerate(layers):
    y = Inches(1.6) + Inches(i * 1.3)
    add_shape(slide, Inches(1.5), y, Inches(10.3), Inches(1.0), color)
    add_text_box(slide, Inches(2.0), y + Inches(0.15), Inches(4), Inches(0.6),
                 layer, font_size=20, font_color=DARK, bold=True)
    add_text_box(slide, Inches(6.5), y + Inches(0.25), Inches(5), Inches(0.5),
                 desc, font_size=16, font_color=GRAY)

# 底部说明
add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.5),
             "💡 所有数据本地化部署，符合医疗数据安全要求",
             font_size=14, font_color=PRIMARY, bold=True)

# ============================================
# Slide 13: 优势总结
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "核心优势总结", "为什么选择口服超声充盈造影 + AI")
add_footer(slide, 13)

advantages = [
    ("🎯 精准", "AI 辅助识别，减少漏诊误诊", "灵敏度 ≥ 85%"),
    ("💰 经济", "成本仅为 CT/MRI 的 1/5", "适合大规模筛查"),
    ("⚡ 高效", "AI 分析 < 3 分钟/例", "大幅提升效率"),
    ("🔒 安全", "无创无辐射，患者接受度高", "可重复检查"),
    ("📊 标准", "统一诊断标准，减少主观差异", "Su-RADS 分级"),
    ("🌐 普惠", "基层医院即可开展", "优质医疗资源下沉"),
]

for i, (icon, title, detail) in enumerate(advantages):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.6) + Inches(row * 2.5)
    
    add_shape(slide, x, y, Inches(3.8), Inches(2.2), LIGHT)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                 icon + " " + title, font_size=18, font_color=PRIMARY, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.0),
                 detail, font_size=14, font_color=GRAY)

# ============================================
# Slide 14: 未来展望
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "未来展望", "从胰腺到多病种的 AI 诊断平台")
add_footer(slide, 14)

# 左侧：短期目标
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.7), Inches(0.5),
             "🔮 短期目标（1-2 年）", font_size=22, font_color=PRIMARY, bold=True)

short_term = [
    "✓ 完成胰腺疾病 AI 模型验证",
    "✓ 积累 1000+ 标注病例",
    "✓ 发表学术论文 2-3 篇",
    "✓ 建立 3-5 家合作医院",
    "✓ 通过伦理审查",
    "✓ 完成 Phase 1 技术验证",
]
add_bullet_list(slide, Inches(1.2), Inches(2.5), Inches(4.7), Inches(4), short_term, font_size=15, spacing=Pt(12))

# 右侧：长期愿景
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
             "🌟 长期愿景（3-5 年）", font_size=22, font_color=SECONDARY, bold=True)

long_term = [
    "🔸 扩展至胃、肝、胆等多病种",
    "🔸 建立国家级超声 AI 数据库",
    "🔸 推动基层医院普及超声筛查",
    "🔸 实现「三甲医院能力下沉」",
    "🔸 申报 NMPA 三类医疗器械",
    "🔸 探索医保覆盖与商业化路径",
]
add_bullet_list(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(4), long_term, font_size=15, spacing=Pt(12))

# ============================================
# Slide 15: Q&A
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# 顶部装饰
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), PRIMARY)

# 主标题
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
             "感谢聆听",
             font_size=48, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
             "Q & A",
             font_size=36, font_color=SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)

# 联系方式
add_shape(slide, Inches(3.5), Inches(4.2), Inches(6.3), Inches(2.0), LIGHT)
add_text_box(slide, Inches(3.5), Inches(4.4), Inches(6.3), Inches(0.5),
             "📧 联系我们", font_size=20, font_color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

contact = [
    "AIMED 充盈视界 · 超声造影人工智能联合实验室",
    "官网：www.aius.xin",
    "地址：浙江省湖州市",
]
add_bullet_list(slide, Inches(3.5), Inches(5.0), Inches(6.3), Inches(1.0), contact, font_size=14, spacing=Pt(6))
# Center the text in the contact list
for p in slide.shapes[-1].text_frame.paragraphs:
    p.alignment = PP_ALIGN.CENTER

# 底部装饰
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), PRIMARY)

# ============================================
# 保存
# ============================================
output_path = "/root/.openclaw/workspace/attachments/口服超声充盈造影胰腺疾病AI诊断模型.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"📊 总页数：{len(prs.slides)}")
